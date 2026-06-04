#!/usr/bin/env python3
"""
Local PPTX operation tool (Layer 3, no backend required).
Dependency: pip install python-pptx

Subcommands:
  info     View basic file info (slide count, slide titles)
  delete   Delete specified slides (1-based, supports multiple and ranges, e.g. 3,5,7-9)
  reorder  Reorder slide sequence (e.g. 2,1,4,3,5)
  extract  Extract selected slides into a new file
  merge    Merge multiple pptx files

Usage examples:
  python scripts/local_pptx_ops.py info --file my.pptx
  python scripts/local_pptx_ops.py delete --file my.pptx --slides 3,5,7-9
  python scripts/local_pptx_ops.py delete --file my.pptx --slides 3,5,7-9 -o trimmed.pptx
  python scripts/local_pptx_ops.py reorder --file my.pptx --order 2,1,4,3,5
  python scripts/local_pptx_ops.py extract --file my.pptx --slides 1-3 -o subset.pptx
  python scripts/local_pptx_ops.py merge --files a.pptx b.pptx -o merged.pptx
"""

import argparse
import os
import sys
import shutil
import copy
import warnings

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from lxml import etree
except ImportError:
    print("Missing dependency, please run: pip install python-pptx", file=sys.stderr)
    sys.exit(1)


# ---------------------------------------------------------------------------
# Parse slide number string (1-based), return 0-based index list
# ---------------------------------------------------------------------------

def parse_slide_spec(spec: str, total: int) -> list[int]:
    """Parse a slide number string like '1,3,5-8,10', return a deduplicated sorted list of 0-based indices."""
    indices = set()
    for part in spec.split(","):
        part = part.strip()
        if not part:
            continue
        if "-" in part:
            lo, hi = part.split("-", 1)
            lo_i, hi_i = int(lo.strip()), int(hi.strip())
            if lo_i < 1 or hi_i > total:
                raise ValueError(f"Slide range {lo_i}-{hi_i} is out of bounds for a file with {total} slides")
            for i in range(lo_i - 1, hi_i):
                indices.add(i)
        else:
            n = int(part)
            if n < 1 or n > total:
                raise ValueError(f"Slide number {n} is out of bounds for a file with {total} slides")
            indices.add(n - 1)
    return sorted(indices)


def parse_order_spec(spec: str, total: int) -> list[int]:
    """Parse an order string for reorder, preserving user-specified order, return 0-based index list."""
    result = []
    seen = set()
    for part in spec.split(","):
        part = part.strip()
        if not part:
            continue
        n = int(part)
        if n < 1 or n > total:
            raise ValueError(f"Slide number {n} is out of bounds for a file with {total} slides")
        idx = n - 1
        if idx in seen:
            raise ValueError(f"Slide number {n} appears more than once")
        seen.add(idx)
        result.append(idx)
    return result


# ---------------------------------------------------------------------------
# Core operation: remove a single slide (XML manipulation, not natively supported by python-pptx)
# ---------------------------------------------------------------------------

def _remove_slide(prs: "Presentation", index: int) -> None:
    """Remove the slide at the given index (0-based) from the Presentation."""
    xml_slides = prs.slides._sldIdLst
    slide_elem = xml_slides[index]
    rId = slide_elem.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
    xml_slides.remove(slide_elem)
    # Remove the corresponding relationship from rels to avoid "Duplicate name" warnings on save
    prs.part.rels.pop(rId)


# ---------------------------------------------------------------------------
# Subcommand: info
# ---------------------------------------------------------------------------

def cmd_info(args):
    prs = Presentation(args.file)
    total = len(prs.slides)
    print(f"File: {os.path.abspath(args.file)}")
    print(f"Total slides: {total}")
    print(f"Slide dimensions: {prs.slide_width.inches:.1f}\" x {prs.slide_height.inches:.1f}\"")
    print()
    for i, slide in enumerate(prs.slides):
        title_text = ""
        for shape in slide.shapes:
            if shape.has_text_frame and shape.shape_type == 13:
                continue  # Skip images
            try:
                if hasattr(shape, "placeholder_format") and shape.placeholder_format is not None:
                    ph_idx = shape.placeholder_format.idx
                    if ph_idx == 0:  # Title placeholder
                        title_text = shape.text_frame.text.strip()
                        break
            except Exception:
                pass
        # Fallback: use the first shape that has text
        if not title_text:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    t = shape.text_frame.text.strip()
                    if t:
                        title_text = t[:60]
                        break
        print(f"  Slide {i+1:2d}: {title_text or '(no title)'}")


# ---------------------------------------------------------------------------
# Subcommand: delete
# ---------------------------------------------------------------------------

def cmd_delete(args):
    prs = Presentation(args.file)
    total = len(prs.slides)
    indices = parse_slide_spec(args.slides, total)
    if not indices:
        print("No valid slide numbers specified, no action taken.")
        return

    # Delete from back to front to avoid index shifting
    for idx in reversed(indices):
        _remove_slide(prs, idx)

    pages_str = ", ".join(str(i + 1) for i in indices)
    remaining = total - len(indices)
    out = _resolve_output(args, args.file)
    prs.save(out)
    print(f"Deleted slide(s) {pages_str} ({len(indices)} deleted, {remaining} remaining)")
    print(f"Saved to: {os.path.abspath(out)}")
    print()
    print(f"RESULT: success")
    print(f"OUTPUT_FILE: {os.path.abspath(out)}")


# ---------------------------------------------------------------------------
# Subcommand: reorder
# ---------------------------------------------------------------------------

def cmd_reorder(args):
    prs = Presentation(args.file)
    total = len(prs.slides)
    order = parse_order_spec(args.order, total)

    if len(order) != total:
        raise ValueError(
            f"--order specified {len(order)} slide numbers, but the file has {total} slides; all slides must be included without duplicates."
        )
    if sorted(order) != list(range(total)):
        raise ValueError("--order must include each slide exactly once (no duplicates, no omissions).")

    xml_slides = prs.slides._sldIdLst
    # Retrieve all sldId elements
    sld_id_elems = list(xml_slides)
    # Rearrange in the new order
    for elem in sld_id_elems:
        xml_slides.remove(elem)
    for new_pos in order:
        xml_slides.append(sld_id_elems[new_pos])

    out = _resolve_output(args, args.file)
    prs.save(out)
    order_str = ",".join(str(i + 1) for i in order)
    print(f"Reordered to [{order_str}], total {total} slides")
    print(f"Saved to: {os.path.abspath(out)}")
    print()
    print(f"RESULT: success")
    print(f"OUTPUT_FILE: {os.path.abspath(out)}")


# ---------------------------------------------------------------------------
# Subcommand: extract
# ---------------------------------------------------------------------------

def cmd_extract(args):
    prs = Presentation(args.file)
    total = len(prs.slides)
    keep = parse_slide_spec(args.slides, total)
    if not keep:
        print("No valid slide numbers specified, no action taken.")
        return

    # Delete slides not in the keep list
    delete_indices = [i for i in range(total) if i not in set(keep)]
    for idx in reversed(delete_indices):
        _remove_slide(prs, idx)

    out = args.output or "extracted.pptx"
    prs.save(out)
    pages_str = ", ".join(str(i + 1) for i in keep)
    print(f"Extracted slide(s) {pages_str} ({len(keep)} slides)")
    print(f"Saved to: {os.path.abspath(out)}")
    print()
    print(f"RESULT: success")
    print(f"OUTPUT_FILE: {os.path.abspath(out)}")


# ---------------------------------------------------------------------------
# Subcommand: merge
# ---------------------------------------------------------------------------

def cmd_merge(args):
    if not args.files or len(args.files) < 2:
        raise ValueError("At least two --files must be provided.")

    # Use the first file as the base
    base_prs = Presentation(args.files[0])
    slide_width = base_prs.slide_width
    slide_height = base_prs.slide_height

    for src_path in args.files[1:]:
        src_prs = Presentation(src_path)
        for slide in src_prs.slides:
            # Add a new slide with a blank layout
            slide_layout = base_prs.slide_layouts[-1]  # Use the last layout (usually blank)
            new_slide = base_prs.slides.add_slide(slide_layout)
            # Copy the spTree (shape tree) contents from the source slide
            new_sp_tree = new_slide.shapes._spTree
            old_sp_tree = slide.shapes._spTree
            # Remove default placeholders from the new slide
            for child in list(new_sp_tree):
                new_sp_tree.remove(child)
            # Copy all child elements
            for child in old_sp_tree:
                new_sp_tree.append(copy.deepcopy(child))

    out = args.output or "merged.pptx"
    base_prs.save(out)
    total = len(base_prs.slides)
    files_str = ", ".join(args.files)
    print(f"Merged: {files_str}")
    print(f"Total {total} slides, saved to: {os.path.abspath(out)}")
    print()
    print(f"RESULT: success")
    print(f"OUTPUT_FILE: {os.path.abspath(out)}")


# ---------------------------------------------------------------------------
# Utility functions
# ---------------------------------------------------------------------------

def _resolve_output(args, original_file: str) -> str:
    """Use -o output path if specified by the user, otherwise overwrite the original file."""
    if hasattr(args, "output") and args.output:
        return args.output
    return original_file


# ---------------------------------------------------------------------------
# CLI entry point
# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(
        description="Local PPTX operation tool (delete/reorder/extract/merge slides)",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  python scripts/local_pptx_ops.py info --file my.pptx
  python scripts/local_pptx_ops.py delete --file my.pptx --slides 3,5,7-9
  python scripts/local_pptx_ops.py delete --file my.pptx --slides 3,5,7-9 -o trimmed.pptx
  python scripts/local_pptx_ops.py reorder --file my.pptx --order 2,1,4,3,5
  python scripts/local_pptx_ops.py extract --file my.pptx --slides 1-3 -o subset.pptx
  python scripts/local_pptx_ops.py merge --files a.pptx b.pptx -o merged.pptx
""",
    )
    sub = parser.add_subparsers(dest="cmd", required=True)

    # info
    p_info = sub.add_parser("info", help="View file info (slide count, titles)")
    p_info.add_argument("--file", required=True, help="Input .pptx file path")

    # delete
    p_del = sub.add_parser("delete", help="Delete specified slides (1-based, supports '3,5,7-9')")
    p_del.add_argument("--file", required=True, help="Input .pptx file path")
    p_del.add_argument("--slides", required=True, help="Slide numbers to delete, e.g. '3,5,7-9'")
    p_del.add_argument("-o", "--output", default="", help="Output path (defaults to overwriting the original file)")

    # reorder
    p_reorder = sub.add_parser("reorder", help="Reorder slide sequence")
    p_reorder.add_argument("--file", required=True, help="Input .pptx file path")
    p_reorder.add_argument("--order", required=True, help="New order, e.g. '2,1,4,3,5' (must include all slides)")
    p_reorder.add_argument("-o", "--output", default="", help="Output path (defaults to overwriting the original file)")

    # extract
    p_extract = sub.add_parser("extract", help="Extract selected slides into a new file")
    p_extract.add_argument("--file", required=True, help="Input .pptx file path")
    p_extract.add_argument("--slides", required=True, help="Slide numbers to keep, e.g. '1-3,5'")
    p_extract.add_argument("-o", "--output", default="", help="Output path (defaults to extracted.pptx)")

    # merge
    p_merge = sub.add_parser("merge", help="Merge multiple pptx files")
    p_merge.add_argument("--files", nargs="+", required=True, help="Input file list (merged in order)")
    p_merge.add_argument("-o", "--output", default="", help="Output path (defaults to merged.pptx)")

    args = parser.parse_args()

    try:
        with warnings.catch_warnings():
            warnings.filterwarnings("ignore", message="Duplicate name", category=UserWarning)
            if args.cmd == "info":
                cmd_info(args)
            elif args.cmd == "delete":
                cmd_delete(args)
            elif args.cmd == "reorder":
                cmd_reorder(args)
            elif args.cmd == "extract":
                cmd_extract(args)
            elif args.cmd == "merge":
                cmd_merge(args)
    except (ValueError, FileNotFoundError) as e:
        print(f"Error: {e}", file=sys.stderr)
        print()
        print(f"RESULT: error")
        print(f"ERROR: {e}")
        sys.exit(1)


if __name__ == "__main__":
    main()
