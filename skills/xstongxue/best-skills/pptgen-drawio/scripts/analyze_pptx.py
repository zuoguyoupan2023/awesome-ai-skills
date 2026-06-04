#!/usr/bin/env python3
"""
analyze_pptx.py - PPT 模板样式提取工具

用法:
  python scripts/analyze_pptx.py ppt_template/模板.pptx reference/style-custom.md
  python scripts/analyze_pptx.py ppt_template/模板.pptx   # 打印到控制台

功能：
  - 提取模板的配色、字体、字号、版式规则
  - 输出可直接用于 pptgen-drawio SKILL 的 reference markdown 文件
  - 若未指定 output.md，则打印到控制台

目录约定（pptgen-drawio skill）：
  - ppt_template/  : 存放参考 PPT 模板
  - scripts/       : 本脚本所在目录
  - reference/     : 样式提取结果输出目录
"""

import sys
import os
sys.stdout.reconfigure(encoding='utf-8')

try:
    from pptx import Presentation
except ImportError:
    print("错误：未安装 python-pptx，请先执行: pip install python-pptx")
    sys.exit(1)

# ─────────────────────────────────────────────
# 辅助函数
# ─────────────────────────────────────────────

def rgb_hex(color_obj):
    """尝试取 RGB 十六进制，失败返回 None"""
    try:
        if color_obj.type is not None:
            return f"#{color_obj.rgb}"
    except Exception:
        pass
    return None


def analyze(path: str) -> dict:
    prs = Presentation(path)
    name = os.path.basename(path)

    slide_w_mm = round(prs.slide_width  / 914400 * 25.4)
    slide_h_mm = round(prs.slide_height / 914400 * 25.4)

    layouts = []
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if layout.name not in layouts:
                layouts.append(layout.name)

    fill_colors  = {}   # color -> count
    font_names   = {}   # name  -> count
    font_sizes   = []
    font_colors  = {}   # color -> count
    slide_samples = []  # [(page, shapes_summary)]

    slides = list(prs.slides)
    total  = len(slides)

    for sidx, slide in enumerate(slides):
        page_info = {"page": sidx + 1, "fills": [], "texts": []}
        for shape in slide.shapes:
            # 填充色
            try:
                fill = shape.fill
                if fill.type == 1:
                    c = rgb_hex(fill.fore_color)
                    if c:
                        fill_colors[c] = fill_colors.get(c, 0) + 1
                        page_info["fills"].append({
                            "color": c,
                            "x_mm": round(shape.left  / 914400 * 25.4),
                            "y_mm": round(shape.top   / 914400 * 25.4),
                            "w_mm": round(shape.width / 914400 * 25.4),
                            "h_mm": round(shape.height/ 914400 * 25.4),
                        })
            except Exception:
                pass

            # 文字
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        text = run.text.strip()
                        if not text:
                            continue
                        font = run.font
                        fname = font.name
                        fsize = round(font.size.pt) if font.size else None
                        fbold = font.bold
                        fcolor = rgb_hex(font.color) if font.color else None

                        if fname:
                            font_names[fname] = font_names.get(fname, 0) + 1
                        if fsize:
                            font_sizes.append(fsize)
                        if fcolor:
                            font_colors[fcolor] = font_colors.get(fcolor, 0) + 1

                        page_info["texts"].append({
                            "text":  text[:40],
                            "font":  fname or "继承",
                            "size":  fsize,
                            "bold":  fbold,
                            "color": fcolor,
                        })

        if sidx < 8:   # 只保存前8页详情
            slide_samples.append(page_info)

    # 排序统计
    top_fills  = sorted(fill_colors.items(),  key=lambda x: -x[1])
    top_fonts  = sorted(font_names.items(),   key=lambda x: -x[1])
    top_fcolors= sorted(font_colors.items(),  key=lambda x: -x[1])

    size_min = min(font_sizes) if font_sizes else None
    size_max = max(font_sizes) if font_sizes else None
    size_avg = round(sum(font_sizes) / len(font_sizes), 1) if font_sizes else None

    return {
        "name":          name,
        "total_slides":  total,
        "slide_w_mm":    slide_w_mm,
        "slide_h_mm":    slide_h_mm,
        "layouts":       layouts,
        "top_fills":     top_fills[:8],
        "top_fonts":     top_fonts[:8],
        "top_fcolors":   top_fcolors[:8],
        "size_min":      size_min,
        "size_max":      size_max,
        "size_avg":      size_avg,
        "slide_samples": slide_samples,
    }


def render_md(data: dict) -> str:
    d = data
    lines = []
    lines.append(f"# PPT 模板样式提取：{d['name']}\n")
    lines.append(f"> 源文件：`{d['name']}`  幻灯片数：{d['total_slides']}  画布：{d['slide_w_mm']} × {d['slide_h_mm']} mm\n")

    lines.append("## 可用版式 (Layouts)\n")
    for l in d["layouts"]:
        lines.append(f"- {l}")
    lines.append("")

    lines.append("## 配色统计（形状填充色，按出现频次排序）\n")
    lines.append("| 色值 | 出现次数 |")
    lines.append("|------|---------|")
    for color, cnt in d["top_fills"]:
        lines.append(f"| `{color}` | {cnt} |")
    lines.append("")

    lines.append("## 字体统计（按出现频次排序）\n")
    lines.append("| 字体名称 | 出现次数 |")
    lines.append("|---------|---------|")
    for fname, cnt in d["top_fonts"]:
        lines.append(f"| `{fname}` | {cnt} |")
    lines.append("")

    lines.append("## 字色统计（按出现频次排序）\n")
    lines.append("| 色值 | 出现次数 |")
    lines.append("|------|---------|")
    for color, cnt in d["top_fcolors"]:
        lines.append(f"| `{color}` | {cnt} |")
    lines.append("")

    lines.append("## 字号范围\n")
    lines.append(f"- 最小：{d['size_min']} pt")
    lines.append(f"- 最大：{d['size_max']} pt")
    lines.append(f"- 平均：{d['size_avg']} pt")
    lines.append("")

    lines.append("## 分页详细分析（前 8 页）\n")
    for page in d["slide_samples"]:
        lines.append(f"### 第 {page['page']} 页\n")
        if page["fills"]:
            lines.append("**形状填充色：**\n")
            for f in page["fills"]:
                lines.append(f"- `{f['color']}` 位置=({f['x_mm']}mm, {f['y_mm']}mm) 宽={f['w_mm']}mm 高={f['h_mm']}mm")
            lines.append("")
        if page["texts"]:
            lines.append("**文字：**\n")
            lines.append("| 内容 | 字体 | 字号 | 粗体 | 颜色 |")
            lines.append("|------|------|------|------|------|")
            for t in page["texts"]:
                lines.append(f"| `{t['text']}` | {t['font']} | {t['size'] or '继承'} | {t['bold']} | {t['color'] or '继承'} |")
            lines.append("")

    lines.append("---\n")
    lines.append("## 使用说明\n")
    lines.append("本文件由 `scripts/analyze_pptx.py` 从 `ppt_template/` 中的模板提取生成。")
    lines.append("放入 `reference/` 目录并在 SKILL.md 的风格表中注册后，即可通过「使用自定义模板」生成对应风格的 Draw.io PPT。")

    return "\n".join(lines)


# ─────────────────────────────────────────────
# 入口
# ─────────────────────────────────────────────

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("用法: python scripts/analyze_pptx.py ppt_template/<模板.pptx> [reference/output.md]")
        sys.exit(1)

    pptx_path = sys.argv[1]
    if not os.path.exists(pptx_path):
        print(f"错误：找不到文件 {pptx_path}")
        sys.exit(1)

    print(f"正在分析 {pptx_path} ...", file=sys.stderr)
    data = analyze(pptx_path)
    md   = render_md(data)

    if len(sys.argv) >= 3:
        out_path = sys.argv[2]
        with open(out_path, "w", encoding="utf-8") as f:
            f.write(md)
        print(f"✅ 样式提取完成，已保存至 {out_path}", file=sys.stderr)
    else:
        print(md)